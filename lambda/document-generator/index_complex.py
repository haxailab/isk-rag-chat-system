import json
import boto3
import logging
import os
import io
import uuid
from datetime import datetime, timedelta
from typing import Dict, Any, Optional
from urllib.parse import quote

# Document generation libraries
try:
    from docx import Document
    from docx.shared import Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.style import WD_STYLE_TYPE
    import docx.oxml
except ImportError:
    logging.warning("python-docx not available - Word generation disabled")

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib import colors
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.units import inch
except ImportError:
    logging.warning("reportlab not available - PDF generation disabled")

try:
    from pptx import Presentation
    from pptx.util import Inches as PPTXInches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as PPTXRGBColor
except ImportError:
    logging.warning("python-pptx not available - PowerPoint generation disabled")

logger = logging.getLogger()
logger.setLevel(logging.INFO)

# ISK Brand Colors
ISK_RED = "#C8102E"
ISK_RED_RGB = (200, 16, 46)
ISK_DARK_GRAY = "#333333"
ISK_LIGHT_GRAY = "#F5F5F5"

def handler(event: Dict[str, Any], context: Any) -> Dict[str, Any]:
    """
    Lambda function to generate documents (Word, PDF, PowerPoint) from RAG responses
    """
    headers = {
        'Content-Type': 'application/json',
        'Access-Control-Allow-Origin': '*',
        'Access-Control-Allow-Headers': 'Content-Type,Authorization,X-Amz-Date,X-Api-Key,X-Amz-Security-Token',
        'Access-Control-Allow-Methods': 'POST,OPTIONS',
        'Access-Control-Allow-Credentials': 'false'
    }

    try:
        logger.info(f"Document generator started. Request ID: {context.aws_request_id}")
        logger.info(f"Event: {json.dumps(event, default=str, ensure_ascii=False)}")

        # Handle preflight request
        if event.get('httpMethod') == 'OPTIONS':
            return {
                'statusCode': 200,
                'headers': headers,
                'body': json.dumps({})
            }

        # Parse request body
        body = json.loads(event['body'])
        content = body.get('content', '')
        title = body.get('title', '無機RAG資料')
        format_type = body.get('format', 'word').lower()  # word, pdf, powerpoint
        sources = body.get('sources', [])
        source_links = body.get('source_links', [])

        if not content:
            return {
                'statusCode': 400,
                'headers': headers,
                'body': json.dumps({'error': 'Content is required'})
            }

        logger.info(f"Generating {format_type} document: {title}")

        # Generate document based on format
        if format_type == 'word':
            file_buffer, filename = generate_word_document(content, title, sources, source_links)
        elif format_type == 'pdf':
            file_buffer, filename = generate_pdf_document(content, title, sources, source_links)
        elif format_type == 'powerpoint':
            file_buffer, filename = generate_powerpoint_document(content, title, sources, source_links)
        else:
            return {
                'statusCode': 400,
                'headers': headers,
                'body': json.dumps({'error': f'Unsupported format: {format_type}'})
            }

        # Upload to S3
        s3_client = boto3.client('s3')
        bucket_name = os.environ.get('TEMP_BUCKET_NAME')

        if not bucket_name:
            logger.error("TEMP_BUCKET_NAME environment variable not set")
            return {
                'statusCode': 500,
                'headers': headers,
                'body': json.dumps({'error': 'Server configuration error'})
            }

        # Generate unique filename
        unique_id = str(uuid.uuid4())[:8]
        s3_key = f"documents/{unique_id}_{filename}"

        # Upload to S3
        s3_client.upload_fileobj(
            file_buffer,
            bucket_name,
            s3_key,
            ExtraArgs={
                'ContentType': get_content_type(format_type),
                'ContentDisposition': f'attachment; filename="{filename}"'
            }
        )

        # Generate presigned URL (15 minutes expiration)
        download_url = s3_client.generate_presigned_url(
            'get_object',
            Params={'Bucket': bucket_name, 'Key': s3_key},
            ExpiresIn=900  # 15 minutes
        )

        logger.info(f"Document generated and uploaded: {s3_key}")

        return {
            'statusCode': 200,
            'headers': headers,
            'body': json.dumps({
                'message': f'{format_type.title()}文書を生成しました',
                'download_url': download_url,
                'filename': filename,
                'format': format_type,
                'expires_in': '15分',
                'file_size': len(file_buffer.getvalue())
            }, ensure_ascii=False)
        }

    except Exception as e:
        logger.error(f"Error in document generator: {str(e)}")
        import traceback
        logger.error(f"Traceback: {traceback.format_exc()}")

        return {
            'statusCode': 500,
            'headers': headers,
            'body': json.dumps({
                'error': '文書生成でエラーが発生しました',
                'details': str(e)
            }, ensure_ascii=False)
        }


def generate_word_document(content: str, title: str, sources: list, source_links: list) -> tuple[io.BytesIO, str]:
    """Generate Word document"""
    doc = Document()

    # Add ISK logo placeholder and title
    title_para = doc.add_heading(title, 0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add subtitle
    subtitle = doc.add_heading('ISK社内資料', level=2)
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Add generation date
    date_para = doc.add_paragraph(f'生成日: {datetime.now().strftime("%Y年%m月%d日 %H:%M")}')
    date_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # Add separator
    doc.add_paragraph('─' * 50)

    # Add main content
    content_heading = doc.add_heading('内容', level=1)
    content_para = doc.add_paragraph(content)

    # Add sources if available
    if sources:
        doc.add_heading('参考文書', level=1)
        for i, source in enumerate(sources, 1):
            doc.add_paragraph(f'{i}. {source}', style='List Number')

    # Add footer
    doc.add_paragraph('─' * 50)
    footer = doc.add_paragraph('無機RAG powered by Claude Sonnet 4.6')
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Save to buffer
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    filename = f"{title}_{datetime.now().strftime('%Y%m%d_%H%M')}.docx"
    return buffer, filename


def generate_pdf_document(content: str, title: str, sources: list, source_links: list) -> tuple[io.BytesIO, str]:
    """Generate PDF document"""
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4,
                           rightMargin=72, leftMargin=72,
                           topMargin=72, bottomMargin=18)

    # Build styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=18,
        spaceAfter=30,
        alignment=1,  # Center
        textColor=colors.HexColor(ISK_RED)
    )

    subtitle_style = ParagraphStyle(
        'CustomSubtitle',
        parent=styles['Heading2'],
        fontSize=14,
        spaceAfter=20,
        alignment=1,
        textColor=colors.HexColor(ISK_DARK_GRAY)
    )

    content_style = ParagraphStyle(
        'ContentStyle',
        parent=styles['Normal'],
        fontSize=11,
        spaceAfter=12,
        leading=16
    )

    # Build document
    story = []

    # Title
    story.append(Paragraph(title, title_style))
    story.append(Paragraph('ISK社内資料', subtitle_style))
    story.append(Spacer(1, 12))

    # Date
    story.append(Paragraph(f'生成日: {datetime.now().strftime("%Y年%m月%d日 %H:%M")}', styles['Normal']))
    story.append(Spacer(1, 20))

    # Content
    story.append(Paragraph('内容', styles['Heading2']))
    story.append(Spacer(1, 12))

    # Split content into paragraphs
    for para in content.split('\n\n'):
        if para.strip():
            story.append(Paragraph(para.replace('\n', '<br/>'), content_style))
            story.append(Spacer(1, 6))

    # Sources
    if sources:
        story.append(Spacer(1, 20))
        story.append(Paragraph('参考文書', styles['Heading2']))
        story.append(Spacer(1, 12))
        for i, source in enumerate(sources, 1):
            story.append(Paragraph(f'{i}. {source}', styles['Normal']))

    # Footer
    story.append(Spacer(1, 30))
    story.append(Paragraph('無機RAG powered by Claude Sonnet 4.6', styles['Normal']))

    # Build PDF
    doc.build(story)
    buffer.seek(0)

    filename = f"{title}_{datetime.now().strftime('%Y%m%d_%H%M')}.pdf"
    return buffer, filename


def generate_powerpoint_document(content: str, title: str, sources: list, source_links: list) -> tuple[io.BytesIO, str]:
    """Generate PowerPoint presentation"""
    prs = Presentation()

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title_placeholder = slide.shapes.title
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = f'ISK社内資料\n{datetime.now().strftime("%Y年%m月%d日")}'

    # Content slide
    content_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(content_slide_layout)

    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = "内容"

    # Split content into bullet points
    content_lines = [line.strip() for line in content.split('\n') if line.strip()]
    text_frame = content_placeholder.text_frame
    text_frame.clear()

    for i, line in enumerate(content_lines[:10]):  # Limit to 10 lines per slide
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = line
        p.level = 0

    # Sources slide
    if sources:
        sources_slide = prs.slides.add_slide(content_slide_layout)
        sources_slide.shapes.title.text = "参考文書"

        sources_text_frame = sources_slide.placeholders[1].text_frame
        sources_text_frame.clear()

        for i, source in enumerate(sources[:10]):  # Limit to 10 sources
            p = sources_text_frame.add_paragraph() if i > 0 else sources_text_frame.paragraphs[0]
            p.text = source
            p.level = 0

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    filename = f"{title}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    return buffer, filename


def get_content_type(format_type: str) -> str:
    """Get MIME type for document format"""
    content_types = {
        'word': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'pdf': 'application/pdf',
        'powerpoint': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    }
    return content_types.get(format_type, 'application/octet-stream')